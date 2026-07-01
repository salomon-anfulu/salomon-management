#!/usr/bin/env python3
"""Generate June 2026 Monthly Report PPT with REAL dynamic scores."""

import json
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ===================== Load real scores =====================
script_dir = os.path.dirname(os.path.abspath(__file__))
with open(os.path.join(script_dir, 'real_scores.json'), 'r', encoding='utf-8') as f:
    all_scores = json.load(f)

# ===================== Brand Colors =====================
SALOMON_DARK = RGBColor(0x0A, 0x1F, 0x44)
SALOMON_RED = RGBColor(0xE2, 0x3B, 0x3B)
SALOMON_BLUE = RGBColor(0x3B, 0x82, 0xF6)
GOLD = RGBColor(0xF5, 0x9E, 0x0B)
SILVER = RGBColor(0x94, 0xA3, 0xB8)
BRONZE = RGBColor(0xCD, 0x7F, 0x32)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xF1, 0xF5, 0xF9)
MID_GRAY = RGBColor(0x64, 0x74, 0x8B)
DARK_TEXT = RGBColor(0x1E, 0x29, 0x3B)
GREEN = RGBColor(0x10, 0xB9, 0x81)
RED = RGBColor(0xEF, 0x44, 0x44)
LIGHT_BLUE = RGBColor(0xDB, 0xEA, 0xFE)
LIGHT_GREEN = RGBColor(0xD1, 0xFA, 0xE5)
LIGHT_RED = RGBColor(0xFE, 0xE2, 0xE2)
LIGHT_GOLD = RGBColor(0xFE, 0xF3, 0xC7)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
W = prs.slide_width
H = prs.slide_height

# ===================== Helpers =====================
def add_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_textbox(slide, left, top, width, height, text, font_size=14, color=DARK_TEXT,
                bold=False, alignment=PP_ALIGN.LEFT, font_name='Microsoft YaHei'):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox

def add_rect(slide, left, top, width, height, fill_color, line_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape

def add_rounded_rect(slide, left, top, width, height, fill_color, line_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape

def add_circle(slide, left, top, diameter, fill_color):
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, diameter, diameter)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape

def add_score_bar(slide, x, y, score, max_score=5, bar_width=Inches(1.5), color=SALOMON_BLUE):
    """Add a visual score bar."""
    add_rect(slide, x, y, bar_width, Inches(0.15), RGBColor(0xE2, 0xE8, 0xF0))
    fill_w = int(bar_width * score / max_score)
    add_rect(slide, x, y, fill_w, Inches(0.15), color)

# ===================== Slide 1: Cover =====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, SALOMON_DARK)
add_rect(slide, Inches(0), Inches(0), Inches(0.15), H, SALOMON_RED)
add_textbox(slide, Inches(1.5), Inches(2.0), Inches(10), Inches(1.0),
            "Salomon 安福路旗舰店", font_size=36, color=WHITE, bold=True)
add_textbox(slide, Inches(1.5), Inches(2.9), Inches(10), Inches(0.8),
            "兼职团队 6月 月度总结", font_size=44, color=WHITE, bold=True)
add_textbox(slide, Inches(1.5), Inches(4.2), Inches(8), Inches(0.5),
            "2026年6月 · Service Team 月度复盘与表彰", font_size=20, color=RGBColor(0x94, 0xA3, 0xB8))
add_textbox(slide, Inches(1.5), Inches(5.8), Inches(6), Inches(0.4),
            "2026年7月 · 安福路 L140", font_size=16, color=RGBColor(0x64, 0x74, 0x8B))
add_circle(slide, Inches(10.8), Inches(0.8), Inches(1.5), RGBColor(0x1A, 0x35, 0x66))
add_circle(slide, Inches(11.5), Inches(5.5), Inches(1.2), RGBColor(0x1A, 0x35, 0x66))

# ===================== Slide 2: TOC =====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_textbox(slide, Inches(1), Inches(0.5), Inches(8), Inches(0.6),
            "目录", font_size=28, color=SALOMON_DARK, bold=True)
add_rect(slide, Inches(1), Inches(1.1), Inches(2), Pt(3), SALOMON_RED)
toc_items = [
    ("01", "系统简介", "兼职管理系统功能概览"),
    ("02", "6月兼职情况总览", "关键数据与整体表现"),
    ("03", "表现评分 TOP 3 表彰", "综合表现前三名奖励"),
    ("04", "表现评分关注分析", "后三名原因分析与改进方向"),
    ("05", "业绩达标表扬", "时产 & UPT KPI 达成者"),
    ("06", "人员变动", "7月团队调整与新成员"),
]
y = Inches(1.6)
for num, title, desc in toc_items:
    add_rounded_rect(slide, Inches(1), y, Inches(0.6), Inches(0.6), SALOMON_DARK)
    add_textbox(slide, Inches(1), y + Inches(0.08), Inches(0.6), Inches(0.5),
                num, font_size=18, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(1.9), y - Inches(0.02), Inches(9), Inches(0.35),
                title, font_size=18, bold=True, color=SALOMON_DARK)
    add_textbox(slide, Inches(1.9), y + Inches(0.32), Inches(9), Inches(0.3),
                desc, font_size=12, color=MID_GRAY)
    y += Inches(0.85)

# ===================== Slide 3: System Intro =====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_textbox(slide, Inches(1), Inches(0.5), Inches(10), Inches(0.6),
            "01  系统简介", font_size=28, color=SALOMON_DARK, bold=True)
add_rect(slide, Inches(1), Inches(1.1), Inches(2), Pt(3), SALOMON_RED)
add_textbox(slide, Inches(1), Inches(1.4), Inches(11), Inches(0.5),
            "Salomon 安福路兼职团队管理系统 —— 数据驱动的兼职管理一体化平台",
            font_size=15, color=MID_GRAY)
features = [
    ("📊", "工作台", "数据概览 · 快捷操作\n团队分布 · 近期动态", SALOMON_BLUE),
    ("📝", "我的填报", "可上班时间 · 换班登记\n店务支援 · 门迎排班", RGBColor(0x83, 0x60, 0xF0)),
    ("👥", "人员管理", "兼职档案 · 部门筛选\n入职/离职管理", GREEN),
    ("📅", "供班总览", "周达标分析 · 排班总览\n逐日填写 · 备注支持", RGBColor(0xF5, 0x9E, 0x0B)),
    ("✅", "考勤记录", "灵工打卡 · 实时考勤\n迟到/缺卡自动识别", RGBColor(0xEC, 0x48, 0x99)),
    ("⭐", "表现评分", "5维度评分 · 时薪联动\n工时/业绩/行为/考勤/好评", SALOMON_RED),
    ("💰", "业绩数据", "销售排行 · 时产UPT\nKPI达标 · 综合对比", RGBColor(0x06, 0xB6, 0xD4)),
    ("🏆", "顾客好评", "大众点评5星好评\n关键词云 · 原文展示", RGBColor(0x84, 0xCC, 0x16)),
]
x_start = Inches(0.8)
y_start = Inches(2.2)
card_w = Inches(2.7)
card_h = Inches(2.1)
for i, (icon, title, desc, color) in enumerate(features):
    col = i % 4
    row = i // 4
    x = x_start + col * (card_w + Inches(0.25))
    y = y_start + row * (card_h + Inches(0.25))
    add_rounded_rect(slide, x, y, card_w, card_h, LIGHT_GRAY)
    add_rounded_rect(slide, x, y, card_w, Inches(0.06), color)
    add_textbox(slide, x + Inches(0.2), y + Inches(0.2), Inches(0.8), Inches(0.6),
                icon, font_size=28, color=DARK_TEXT)
    add_textbox(slide, x + Inches(0.2), y + Inches(0.8), Inches(2.3), Inches(0.4),
                title, font_size=16, color=SALOMON_DARK, bold=True)
    add_textbox(slide, x + Inches(0.2), y + Inches(1.25), Inches(2.3), Inches(0.8),
                desc, font_size=11, color=MID_GRAY)

# ===================== Slide 4: June Overview =====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_textbox(slide, Inches(1), Inches(0.5), Inches(10), Inches(0.6),
            "02  6月兼职情况总览", font_size=28, color=SALOMON_DARK, bold=True)
add_rect(slide, Inches(1), Inches(1.1), Inches(2), Pt(3), SALOMON_RED)
metrics = [
    ("¥293,250", "6月总销售额", SALOMON_DARK),
    ("1,625.5h", "总出勤工时", SALOMON_BLUE),
    ("13人", "Service Team", GREEN),
    ("1.36", "团队平均 UPT", RGBColor(0x83, 0x60, 0xF0)),
]
x_start = Inches(0.8)
card_w = Inches(2.8)
for i, (val, label, color) in enumerate(metrics):
    x = x_start + i * (card_w + Inches(0.2))
    y = Inches(1.5)
    add_rounded_rect(slide, x, y, card_w, Inches(1.5), color)
    add_textbox(slide, x + Inches(0.1), y + Inches(0.2), card_w - Inches(0.2), Inches(0.7),
                val, font_size=28, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, x + Inches(0.1), y + Inches(0.95), card_w - Inches(0.2), Inches(0.4),
                label, font_size=13, color=WHITE, alignment=PP_ALIGN.CENTER)
add_textbox(slide, Inches(1), Inches(3.3), Inches(5), Inches(0.5),
            "关键运营指标", font_size=16, color=SALOMON_DARK, bold=True)
detail_items = [
    "门迎排班：覆盖整月，日均10+班次，门迎总时长超过300小时",
    "顾客好评：大众点评5星好评10条，涉及4名兼职（迟骋4条/朱凯赟2条/杨子豪2条/陈昕媛1条/李若彤1条）",
    "考勤异常：迟到3人次（迟骋/王靳毓/李若彤×2），旷工1人次（李若彤6/6取消）",
    "供班达标率：多数成员每周供班4天+周末1天达标，王龙宇因出差19-30日供班严重不足",
]
y = Inches(3.9)
for item in detail_items:
    add_rect(slide, Inches(1.2), y + Inches(0.1), Inches(0.1), Inches(0.1), SALOMON_RED)
    add_textbox(slide, Inches(1.5), y, Inches(10.5), Inches(0.5), item, font_size=13, color=DARK_TEXT)
    y += Inches(0.55)
add_textbox(slide, Inches(1), Inches(6.2), Inches(11), Inches(0.5),
            "趋势对比：4月 ¥407,876 → 5月 ¥440,577 → 6月 ¥293,250（6月因人员调整及排班减少，总销售有所下降）",
            font_size=12, color=MID_GRAY)

# ===================== Slide 5: TOP 3 (REAL DATA) =====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_textbox(slide, Inches(1), Inches(0.5), Inches(10), Inches(0.6),
            "03  表现评分 TOP 3 表彰", font_size=28, color=SALOMON_DARK, bold=True)
add_rect(slide, Inches(1), Inches(1.1), Inches(2), Pt(3), GOLD)
add_textbox(slide, Inches(1), Inches(1.3), Inches(11), Inches(0.5),
            "综合表现评分 = 工时支持 + 销售业绩 + 行为规范 + 考勤纪律 + 顾客好评（5维度平均）",
            font_size=13, color=MID_GRAY)

# TOP 3 from real data
top3 = all_scores[:3]
rewards = ["480元", "360元", "240元"]
medals = ["🥇", "🥈", "🥉"]
colors_rank = [GOLD, SILVER, BRONZE]
bg_colors = [LIGHT_GOLD, RGBColor(0xF1, 0xF5, 0xF9), RGBColor(0xFE, 0xF6, 0xEC)]

# Highlights for each
highlights_map = {
    "孔祥宇": [
        "18天出勤127.5h，全勤无异常",
        "销售¥22,959，UPT 1.77（全队并列第一）",
        "服装连带率极高：服装占比35.7%",
        "工时、考勤、业绩三维度均≥4分",
        "门迎20h+店务12.7h，行为排名#3",
    ],
    "迟骋": [
        "15天出勤109.5h，UPT 2.13（全队第一！）",
        "大众点评好评4条（全队最多）",
        "好评维度3.5分，拉高综合评分",
        "门迎16.5h+店务6.5h，积极参与团队工作",
        "建议提升时产（当前¥142.3/h）",
    ],
    "朱凯赟": [
        "14天出勤105.5h，UPT 1.50（优秀）",
        "大众点评好评2条",
        "考勤满分5分，全勤无异常",
        "服装连带能力出色（服装占比39%）",
        "建议提升时产和销售总量",
    ],
}

x_start = Inches(0.6)
card_w = Inches(4.0)
card_h = Inches(5.2)
for i, person in enumerate(top3):
    x = x_start + i * (card_w + Inches(0.15))
    y = Inches(1.9)
    add_rounded_rect(slide, x, y, card_w, card_h, bg_colors[i])
    add_rounded_rect(slide, x, y, card_w, Inches(0.08), colors_rank[i])
    add_textbox(slide, x + Inches(0.2), y + Inches(0.2), Inches(2), Inches(0.5),
                medals[i] + "  第" + str(i+1) + "名", font_size=22, color=colors_rank[i], bold=True)
    add_textbox(slide, x + Inches(0.2), y + Inches(0.8), Inches(2), Inches(0.5),
                person["name"], font_size=24, color=SALOMON_DARK, bold=True)
    add_textbox(slide, x + Inches(2.5), y + Inches(0.8), Inches(1.3), Inches(0.5),
                str(person["avg"]), font_size=32, color=colors_rank[i], bold=True, alignment=PP_ALIGN.RIGHT)
    add_textbox(slide, x + Inches(3.5), y + Inches(0.9), Inches(0.4), Inches(0.3),
                "分", font_size=12, color=MID_GRAY)
    add_rounded_rect(slide, x + Inches(0.2), y + Inches(1.4), card_w - Inches(0.4), Inches(0.5), colors_rank[i])
    add_textbox(slide, x + Inches(0.2), y + Inches(1.47), card_w - Inches(0.4), Inches(0.4),
                "奖金 ¥" + rewards[i], font_size=18, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    # Score breakdown
    y_scores = y + Inches(2.1)
    add_textbox(slide, x + Inches(0.2), y_scores, Inches(3.5), Inches(0.3),
                "维度评分", font_size=11, color=MID_GRAY, bold=True)
    y_scores += Inches(0.3)
    dim_names = {"availability": "工时支持", "performance": "销售业绩", "behavior": "行为规范", "attendance": "考勤纪律", "customerReview": "顾客好评"}
    for dim_key, dim_label in dim_names.items():
        sc = person["scores"][dim_key]
        add_textbox(slide, x + Inches(0.3), y_scores, Inches(1.5), Inches(0.3),
                    dim_label, font_size=11, color=DARK_TEXT)
        add_score_bar(slide, x + Inches(1.8), y_scores + Inches(0.05), sc, color=colors_rank[i])
        add_textbox(slide, x + Inches(3.35), y_scores, Inches(0.4), Inches(0.3),
                    str(sc), font_size=11, color=colors_rank[i], bold=True)
        y_scores += Inches(0.3)
    # Highlights
    y_hl = y + Inches(4.0)
    add_textbox(slide, x + Inches(0.2), y_hl, Inches(3.5), Inches(0.3),
                "表现亮点", font_size=11, color=MID_GRAY, bold=True)
    y_hl += Inches(0.25)
    for hl in highlights_map.get(person["name"], []):
        add_rect(slide, x + Inches(0.25), y_hl + Inches(0.07), Inches(0.06), Inches(0.06), colors_rank[i])
        add_textbox(slide, x + Inches(0.45), y_hl, Inches(3.3), Inches(0.3), hl, font_size=10, color=DARK_TEXT)
        y_hl += Inches(0.25)

# ===================== Slide 6: Bottom 3 (REAL DATA) =====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_textbox(slide, Inches(1), Inches(0.5), Inches(11), Inches(0.6),
            "04  表现评分后三位 · 关注与改进", font_size=28, color=SALOMON_DARK, bold=True)
add_rect(slide, Inches(1), Inches(1.1), Inches(2), Pt(3), RED)
add_textbox(slide, Inches(1), Inches(1.3), Inches(11), Inches(0.4),
            "以下三位综合评分较低，需要重点关注并制定改进计划", font_size=13, color=MID_GRAY)

bottom3 = all_scores[-3:]  # Last 3
issues_map = {
    "邓奇缘": [
        "销售业绩 2.0分（最低项）：时产¥116.4/h，不达标",
        "鞋履占比90.7%，品类过于单一",
        "17天出勤137.5h但销售仅¥16,005",
    ],
    "李若彤": [
        "考勤纪律 1分（最低项）：旷工1次+迟到2次",
        "6/6排班取消=旷工(-2)，6/11+6/30迟到(-2)",
        "其余4维度均≥4分，考勤是唯一短板",
    ],
    "王龙宇": [
        "工时支持 1.5分：4周供班未达标，出勤严重不足",
        "销售业绩 2.0分：时产¥93.2/h，不达标",
        "6月仅出勤8天/58.5h（19-30日出差请假）",
    ],
}
plans_map = {
    "邓奇缘": "建议：拓展服装和配件产品知识，提升连带销售能力；UPT 1.25已达标，时产是主要瓶颈",
    "李若彤": "建议：考勤是唯一短板！加强排班确认，避免取消/迟到。其余维度表现优秀（业绩4分/行为4.7分）",
    "王龙宇": "建议：7月恢复全勤后重点观察供班和时产表现；出差期间供班缺口已由其他成员补位",
}

x_start = Inches(0.6)
card_w = Inches(4.0)
card_h = Inches(5.0)
for i, person in enumerate(bottom3):
    x = x_start + i * (card_w + Inches(0.15))
    y = Inches(1.8)
    add_rounded_rect(slide, x, y, card_w, card_h, RGBColor(0xFF, 0xFA, 0xFA))
    add_rounded_rect(slide, x, y, card_w, Inches(0.08), RED)
    add_textbox(slide, x + Inches(0.2), y + Inches(0.2), Inches(2), Inches(0.5),
                person["name"], font_size=22, color=SALOMON_DARK, bold=True)
    add_textbox(slide, x + Inches(2.8), y + Inches(0.25), Inches(1), Inches(0.4),
                "第" + str(13 - i) + "名", font_size=14, color=RED, bold=True, alignment=PP_ALIGN.RIGHT)
    add_textbox(slide, x + Inches(0.2), y + Inches(0.75), Inches(2), Inches(0.5),
                str(person["avg"]) + " 分", font_size=28, color=RED, bold=True)
    # Score breakdown
    y_s = y + Inches(1.4)
    add_textbox(slide, x + Inches(0.2), y_s, Inches(3.5), Inches(0.3),
                "维度评分", font_size=11, color=MID_GRAY, bold=True)
    y_s += Inches(0.3)
    dim_names = {"availability": "工时支持", "performance": "销售业绩", "behavior": "行为规范", "attendance": "考勤纪律", "customerReview": "顾客好评"}
    for dim_key, dim_label in dim_names.items():
        sc = person["scores"][dim_key]
        bar_color = RED if sc <= 2.5 else GREEN
        add_textbox(slide, x + Inches(0.3), y_s, Inches(1.5), Inches(0.3),
                    dim_label, font_size=11, color=DARK_TEXT)
        add_score_bar(slide, x + Inches(1.8), y_s + Inches(0.05), sc, color=bar_color)
        add_textbox(slide, x + Inches(3.35), y_s, Inches(0.4), Inches(0.3),
                    str(sc), font_size=11, color=bar_color, bold=True)
        y_s += Inches(0.3)
    # Issues
    y_i = y + Inches(3.3)
    add_textbox(slide, x + Inches(0.2), y_i, Inches(3.5), Inches(0.3),
                "问题分析", font_size=11, color=RED, bold=True)
    y_i += Inches(0.28)
    for issue in issues_map.get(person["name"], []):
        add_rect(slide, x + Inches(0.25), y_i + Inches(0.05), Inches(0.06), Inches(0.06), RED)
        add_textbox(slide, x + Inches(0.45), y_i, Inches(3.3), Inches(0.5), issue, font_size=9, color=DARK_TEXT)
        y_i += Inches(0.38)
    # Plan
    y_p = y + Inches(4.5)
    add_rounded_rect(slide, x + Inches(0.15), y_p, card_w - Inches(0.3), Inches(0.45), LIGHT_GREEN)
    add_textbox(slide, x + Inches(0.25), y_p + Inches(0.05), card_w - Inches(0.5), Inches(0.4),
                plans_map.get(person["name"], ""), font_size=9, color=RGBColor(0x06, 0x65, 0x3D))

# ===================== Slide 7: KPI Blind Box =====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_textbox(slide, Inches(1), Inches(0.5), Inches(11), Inches(0.6),
            "05  业绩达标表扬 · 盲盒奖励", font_size=28, color=SALOMON_DARK, bold=True)
add_rect(slide, Inches(1), Inches(1.1), Inches(2), Pt(3), GREEN)
add_rounded_rect(slide, Inches(1), Inches(1.4), Inches(11.3), Inches(0.6), LIGHT_GREEN)
add_textbox(slide, Inches(1.3), Inches(1.5), Inches(10), Inches(0.4),
            "KPI 达标标准：时产 ≥ ¥210/h  且  UPT ≥ 1.25  ——  达成者获得盲盒奖励",
            font_size=14, color=RGBColor(0x06, 0x65, 0x3D), bold=True)

# KPI winners: those with hourly>=210 AND upt>=1.25
kpi_winners = []
for p in all_scores:
    if p["perf"]["hourly"] >= 210 and p["perf"]["upt"] >= 1.25:
        kpi_winners.append(p)

kpi_notes = {
    "杨子豪": "时产全队第一，销售冠军",
    "李若彤": "销量+UPT双达标，稳定输出",
    "王靳毓": "UPT表现优秀，连带能力强",
}

x_start = Inches(1.5)
card_w = Inches(3.3)
card_h = Inches(4.5)
for i, w in enumerate(kpi_winners):
    x = x_start + i * (card_w + Inches(0.3))
    y = Inches(2.3)
    add_rounded_rect(slide, x, y, card_w, card_h, LIGHT_GREEN)
    add_rounded_rect(slide, x, y, card_w, Inches(0.08), GREEN)
    add_textbox(slide, x + Inches(0.2), y + Inches(0.2), Inches(2), Inches(0.5),
                "🎁 盲盒奖励", font_size=14, color=GREEN, bold=True)
    add_textbox(slide, x + Inches(0.2), y + Inches(0.7), Inches(2.5), Inches(0.5),
                w["name"], font_size=24, color=SALOMON_DARK, bold=True)
    y_m = y + Inches(1.4)
    perf = w["perf"]
    metrics_pairs = [
        ("销售额", "¥{:,.0f}".format(perf["sales"])),
        ("时产", "¥{}/h".format(perf["hourly"])),
        ("UPT", str(perf["upt"])),
    ]
    for label, val in metrics_pairs:
        add_textbox(slide, x + Inches(0.2), y_m, Inches(1.2), Inches(0.3),
                    label, font_size=11, color=MID_GRAY)
        add_textbox(slide, x + Inches(1.5), y_m, Inches(1.6), Inches(0.3),
                    val, font_size=12, color=SALOMON_DARK, bold=True)
        y_m += Inches(0.35)
    y_n = y + Inches(3.0)
    add_rounded_rect(slide, x + Inches(0.15), y_n, card_w - Inches(0.3), Inches(0.5), WHITE)
    add_textbox(slide, x + Inches(0.25), y_n + Inches(0.08), card_w - Inches(0.5), Inches(0.4),
                "🏆 " + kpi_notes.get(w["name"], ""), font_size=11, color=GREEN, bold=True)
    reviews_count = w["review"]["count"]
    if reviews_count > 0:
        add_textbox(slide, x + Inches(0.2), y + Inches(3.6), Inches(3), Inches(0.3),
                    "好评 +{} 条".format(reviews_count), font_size=10, color=RGBColor(0x06, 0x65, 0x3D))

add_textbox(slide, Inches(1), Inches(6.9), Inches(11), Inches(0.4),
            "注：6月共13人参与销售统计，仅{}人同时达成时产≥210/h 且 UPT≥1.25 的双KPI标准".format(len(kpi_winners)),
            font_size=11, color=MID_GRAY)

# ===================== Slide 8: Near-miss =====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_textbox(slide, Inches(1), Inches(0.5), Inches(11), Inches(0.6),
            "05  业绩数据全景 · 差一点也值得鼓励", font_size=24, color=SALOMON_DARK, bold=True)
add_rect(slide, Inches(1), Inches(1.1), Inches(2), Pt(3), SALOMON_BLUE)
add_textbox(slide, Inches(1), Inches(1.3), Inches(11), Inches(0.4),
            "以下兼职在时产或UPT单项表现突出，距离双达标仅一步之遥", font_size=13, color=MID_GRAY)

headers = ["姓名", "销售额", "时产(¥/h)", "UPT", "时产状态", "UPT状态", "亮点"]
col_widths = [Inches(1.2), Inches(1.5), Inches(1.4), Inches(1.0), Inches(1.3), Inches(1.3), Inches(3.2)]
x = Inches(0.8)
y_h = Inches(2.0)
add_rect(slide, x, y_h, sum(col_widths, Inches(0)), Inches(0.45), SALOMON_DARK)
cx = x
for i, h in enumerate(headers):
    add_textbox(slide, cx + Inches(0.05), y_h + Inches(0.08), col_widths[i], Inches(0.3),
                h, font_size=12, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    cx += col_widths[i]

# Near-miss: hourly>=180 or upt>=1.25 but NOT both >=210 and >=1.25
near_miss = []
for p in all_scores:
    perf = p["perf"]
    hp = perf["hourly"] >= 210
    up = perf["upt"] >= 1.25
    if not (hp and up) and (perf["hourly"] >= 180 or perf["upt"] >= 1.25):
        note = ""
        if hp and not up:
            note = "时产达标，UPT差{:.2f}即达标".format(1.25 - perf["upt"])
        elif up and not hp:
            if perf["upt"] >= 1.6:
                note = "UPT全队前列，时产需提升"
            else:
                note = "UPT达标，时产差¥{:.0f}即达标".format(210 - perf["hourly"])
        else:
            note = "潜力选手，两个维度均接近达标"
        near_miss.append((p["name"], perf["sales"], perf["hourly"], perf["upt"], hp, up, note))

y_r = y_h + Inches(0.45)
row_h = Inches(0.45)
for idx, (name, sales, hourly, upt, hp, up, note) in enumerate(near_miss):
    bg = LIGHT_GRAY if idx % 2 == 0 else WHITE
    add_rect(slide, x, y_r, sum(col_widths, Inches(0)), row_h, bg)
    cx = x
    cells = [name, "¥{:,.0f}".format(sales), "¥{}".format(hourly), str(upt),
             "✅ 达标" if hp else "❌ 未达标",
             "✅ 达标" if up else "❌ 未达标", note]
    for i, val in enumerate(cells):
        if i == 4: c = GREEN if hp else RED
        elif i == 5: c = GREEN if up else RED
        else: c = DARK_TEXT
        add_textbox(slide, cx + Inches(0.05), y_r + Inches(0.08), col_widths[i], Inches(0.3),
                    val, font_size=11, color=c, bold=(i == 0), alignment=PP_ALIGN.CENTER)
        cx += col_widths[i]
    y_r += row_h

add_textbox(slide, Inches(1), y_r + Inches(0.2), Inches(11), Inches(0.4),
            "分析：时产达标但UPT未达标 → 需提升连带销售能力；UPT达标但时产未达标 → 需提升客单价或增加高峰时段排班",
            font_size=11, color=MID_GRAY)

# ===================== Slide 9: Personnel Changes =====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_textbox(slide, Inches(1), Inches(0.5), Inches(11), Inches(0.6),
            "06  7月人员变动", font_size=28, color=SALOMON_DARK, bold=True)
add_rect(slide, Inches(1), Inches(1.1), Inches(2), Pt(3), SALOMON_BLUE)
add_textbox(slide, Inches(1), Inches(1.3), Inches(11), Inches(0.4),
            "7月起团队结构调整，恭喜、欢迎、期待！", font_size=14, color=MID_GRAY)

changes = [
    {
        "icon": "🎉", "title": "转正晋升", "name": "陈昕媛",
        "tag": "兼职 → 全职", "color": GOLD, "bg": LIGHT_GOLD,
        "desc": "恭喜陈昕媛正式转为全职员工！\n\n6月综合评分3.5分（全队#4），"
                "18天出勤139.5h，销售¥34,360，时产¥246.3/h（KPI达标），"
                "门迎19.5h，大众点评好评1条。\n\n"
                "从1月入职以来持续稳定输出，实力获得认可。",
    },
    {
        "icon": "👋", "title": "新成员加入", "name": "唐蓉",
        "tag": "新加入 Service Team", "color": GREEN, "bg": LIGHT_GREEN,
        "desc": "欢迎唐蓉加入 Service Team！\n\n"
                "作为7月新成员，将参与门店日常服务、销售、门迎等全链路工作。\n\n"
                "期待在新环境中快速成长，融入团队文化。",
    },
    {
        "icon": "🔄", "title": "岗位调整", "name": "玛依拉",
        "tag": "仓库兼职 → Service Team", "color": SALOMON_BLUE, "bg": LIGHT_BLUE,
        "desc": "玛依拉从仓库兼职转入 Service Team！\n\n"
                "凭借在仓库岗位的出色表现和积极态度，"
                "获得转入前场服务团队的机会。\n\n"
                "期待在前场继续发光发热，拓展新的能力边界。",
    },
]

x_start = Inches(0.8)
card_w = Inches(3.8)
card_h = Inches(4.8)
for i, ch in enumerate(changes):
    x = x_start + i * (card_w + Inches(0.2))
    y = Inches(2.0)
    add_rounded_rect(slide, x, y, card_w, card_h, ch["bg"])
    add_rounded_rect(slide, x, y, card_w, Inches(0.08), ch["color"])
    add_textbox(slide, x + Inches(0.2), y + Inches(0.2), Inches(1), Inches(0.6),
                ch["icon"], font_size=32, color=DARK_TEXT)
    add_rounded_rect(slide, x + Inches(0.2), y + Inches(0.85), Inches(2.2), Inches(0.35), ch["color"])
    add_textbox(slide, x + Inches(0.3), y + Inches(0.88), Inches(2), Inches(0.3),
                ch["tag"], font_size=11, color=WHITE, bold=True)
    add_textbox(slide, x + Inches(0.2), y + Inches(1.35), Inches(3), Inches(0.3),
                ch["title"], font_size=14, color=MID_GRAY)
    add_textbox(slide, x + Inches(0.2), y + Inches(1.65), Inches(3), Inches(0.6),
                ch["name"], font_size=28, color=ch["color"], bold=True)
    add_textbox(slide, x + Inches(0.2), y + Inches(2.4), card_w - Inches(0.4), Inches(2.3),
                ch["desc"], font_size=12, color=DARK_TEXT)

# ===================== Slide 10: Closing =====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, SALOMON_DARK)
add_rect(slide, Inches(0), Inches(0), Inches(0.15), H, SALOMON_RED)
add_textbox(slide, Inches(1.5), Inches(2.0), Inches(10), Inches(1.0),
            "感谢每一位的付出", font_size=40, color=WHITE, bold=True)
add_textbox(slide, Inches(1.5), Inches(3.2), Inches(10), Inches(0.6),
            "7月，一起继续冲！", font_size=28, color=GOLD, bold=True)
add_textbox(slide, Inches(1.5), Inches(4.5), Inches(8), Inches(0.5),
            "Salomon 安福路旗舰店 L140 · Service Team", font_size=16, color=RGBColor(0x94, 0xA3, 0xB8))
add_textbox(slide, Inches(1.5), Inches(5.0), Inches(8), Inches(0.5),
            "2026年7月", font_size=14, color=RGBColor(0x64, 0x74, 0x8B))
add_circle(slide, Inches(10.8), Inches(0.8), Inches(1.5), RGBColor(0x1A, 0x35, 0x66))
add_circle(slide, Inches(11.5), Inches(5.5), Inches(1.2), RGBColor(0x1A, 0x35, 0x66))

# ===================== Save =====================
output_path = os.path.join(os.path.dirname(script_dir), "Salomon安福路_6月月度总结.pptx")
prs.save(output_path)
print(f"PPT saved to: {output_path}")
print(f"\nScore ranking used:")
for i, p in enumerate(all_scores):
    print(f"  #{i+1} {p['name']}: {p['avg']} 分")
